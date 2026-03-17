import os
import sys
import json
import re
import copy
import tempfile
from datetime import datetime, timedelta
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
import html as html_mod
import requests


# --- Portable path constants ---
SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
SKILL_DIR = os.path.dirname(SCRIPTS_DIR)
ACTIVE_STYLE_PATH = os.path.join(SKILL_DIR, "active_style.json")
STYLES_DIR = os.path.join(SKILL_DIR, "styles")

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# --- Style resolution ---
def _get_active_style_name():
    """Read active_style.json and return the active style name."""
    with open(ACTIVE_STYLE_PATH, encoding="utf-8") as f:
        return json.load(f)["active"]


def _get_style_dir(style_name=None):
    """Return the directory path for the given style (or the active style)."""
    if style_name is None:
        style_name = _get_active_style_name()
    return os.path.join(STYLES_DIR, style_name)


def _load_style_config(style_name=None):
    """Load style_config.json for the given style (or the active style)."""
    style_dir = _get_style_dir(style_name)
    config_path = os.path.join(style_dir, "style_config.json")
    with open(config_path, encoding="utf-8") as f:
        return json.load(f)


STYLE = _load_style_config()
TEMPLATE_PATH = os.path.join(_get_style_dir(), "template.pptx")

# --- Style constants (loaded from config) ---
LAYOUT_TITLE = STYLE["layouts"]["title"]["name"]
LAYOUT_CONTENT = STYLE["layouts"]["content"]["name"]

PH_TITLE_IDX = STYLE["placeholders"]["title_idx"]
PH_BODY_IDX = STYLE["placeholders"]["body_idx"]
PH_URL_IDX = STYLE["placeholders"]["url_idx"]

FONT_SIZE_BODY = STYLE["fonts"]["body_size_pt100"] * 100 + (STYLE["fonts"]["body_size_pt100"] % 100) * 0
FONT_SIZE_BODY = int(STYLE["fonts"]["body_size_pt100"] * 12700 / 100)
FONT_SIZE_URL = int(STYLE["fonts"]["url_size_pt100"] * 12700 / 100)
FONT_SIZE_TABLE = int(STYLE["fonts"]["table_size_pt100"] * 12700 / 100)
FONT_SIZE_NEXT_MEETING = int(STYLE["fonts"]["next_meeting_size_pt100"] * 12700 / 100)
FONT_SIZE_TITLE_TB = int(STYLE["fonts"]["title_textbox_size_pt"] * 12700)

BULLET_CHAR = STYLE["bullet"]["char"]
BULLET_FONT = STYLE["bullet"]["font"]
BULLET_FONT_PITCH_FAMILY = STYLE["bullet"]["pitch_family"]
BULLET_FONT_CHARSET = STYLE["bullet"]["charset"]
BULLET_SIZE_PCT = STYLE["bullet"]["size_pct"]

BULLET_MARGIN_L0 = STYLE["bullet"]["margin_l0"]
BULLET_MARGIN_L1 = STYLE["bullet"]["margin_l1"]

_tb = STYLE["title_slide_textboxes"]
TITLE_ISSUE_LEFT = _tb["issue"]["left"]
TITLE_ISSUE_TOP = _tb["issue"]["top"]
TITLE_ISSUE_WIDTH = _tb["issue"]["width"]
TITLE_ISSUE_HEIGHT = _tb["issue"]["height"]
TITLE_DATE_LEFT = _tb["date"]["left"]
TITLE_DATE_TOP = _tb["date"]["top"]
TITLE_DATE_WIDTH = _tb["date"]["width"]
TITLE_DATE_HEIGHT = _tb["date"]["height"]
TITLE_TB_FONT = _tb.get("font", "Zen Kaku Gothic New")
TITLE_TB_SIZE_PT = _tb.get("size_pt", 12)
TITLE_TB_BOLD = _tb.get("bold", True)
TITLE_TB_COLOR = _tb.get("color_rgb", "FFFFFF")
TITLE_TB_LINE_SPACING = _tb.get("line_spacing", 1.25)

_tbl = STYLE["table_position"]
TABLE_LEFT = _tbl["left"]
TABLE_TOP = _tbl["top"]
TABLE_WIDTH = _tbl["width"]
TABLE_HEIGHT = _tbl["height"]

_img = STYLE["image_layout"]
BODY_NARROW_WIDTH = _img["body_narrow_width"]
IMG_RIGHT_LEFT = _img["img_right_left"]
IMG_RIGHT_TOP = _img["img_right_top"]
IMG_RIGHT_MAX_WIDTH = _img["img_right_max_width"]
IMG_RIGHT_MAX_HEIGHT = _img["img_right_max_height"]

URL_HYPERLINK_COLOR = STYLE.get("url_hyperlink_color", "0563C1")
URL_SEPARATOR_COLOR = STYLE.get("url_separator_color", "808080")
BODY_BOLD = STYLE.get("body_bold", False)

IMG_TEMP_DIR = os.path.join(SCRIPTS_DIR, ".img_cache")


def reload_style(style_name=None):
    """Reload style configuration.

    Call this after style_analyzer.py has updated the config file
    to pick up new values without restarting.

    Args:
        style_name: specific style to load. If None, reads from active_style.json.
    """
    global STYLE, TEMPLATE_PATH, LAYOUT_TITLE, LAYOUT_CONTENT
    global PH_TITLE_IDX, PH_BODY_IDX, PH_URL_IDX
    global FONT_SIZE_BODY, FONT_SIZE_URL, FONT_SIZE_TABLE
    global FONT_SIZE_NEXT_MEETING, FONT_SIZE_TITLE_TB
    global BULLET_CHAR, BULLET_FONT, BULLET_FONT_PITCH_FAMILY
    global BULLET_FONT_CHARSET, BULLET_SIZE_PCT
    global BULLET_MARGIN_L0, BULLET_MARGIN_L1
    global TITLE_ISSUE_LEFT, TITLE_ISSUE_TOP, TITLE_ISSUE_WIDTH, TITLE_ISSUE_HEIGHT
    global TITLE_DATE_LEFT, TITLE_DATE_TOP, TITLE_DATE_WIDTH, TITLE_DATE_HEIGHT
    global TITLE_TB_FONT, TITLE_TB_SIZE_PT, TITLE_TB_BOLD, TITLE_TB_COLOR, TITLE_TB_LINE_SPACING
    global TABLE_LEFT, TABLE_TOP, TABLE_WIDTH, TABLE_HEIGHT
    global BODY_NARROW_WIDTH, IMG_RIGHT_LEFT, IMG_RIGHT_TOP
    global IMG_RIGHT_MAX_WIDTH, IMG_RIGHT_MAX_HEIGHT
    global URL_HYPERLINK_COLOR, URL_SEPARATOR_COLOR, BODY_BOLD

    STYLE = _load_style_config(style_name)
    TEMPLATE_PATH = os.path.join(_get_style_dir(style_name), "template.pptx")
    LAYOUT_TITLE = STYLE["layouts"]["title"]["name"]
    LAYOUT_CONTENT = STYLE["layouts"]["content"]["name"]
    PH_TITLE_IDX = STYLE["placeholders"]["title_idx"]
    PH_BODY_IDX = STYLE["placeholders"]["body_idx"]
    PH_URL_IDX = STYLE["placeholders"]["url_idx"]
    FONT_SIZE_BODY = int(STYLE["fonts"]["body_size_pt100"] * 12700 / 100)
    FONT_SIZE_URL = int(STYLE["fonts"]["url_size_pt100"] * 12700 / 100)
    FONT_SIZE_TABLE = int(STYLE["fonts"]["table_size_pt100"] * 12700 / 100)
    FONT_SIZE_NEXT_MEETING = int(STYLE["fonts"]["next_meeting_size_pt100"] * 12700 / 100)
    FONT_SIZE_TITLE_TB = int(STYLE["fonts"]["title_textbox_size_pt"] * 12700)
    BULLET_CHAR = STYLE["bullet"]["char"]
    BULLET_FONT = STYLE["bullet"]["font"]
    BULLET_FONT_PITCH_FAMILY = STYLE["bullet"]["pitch_family"]
    BULLET_FONT_CHARSET = STYLE["bullet"]["charset"]
    BULLET_SIZE_PCT = STYLE["bullet"]["size_pct"]
    BULLET_MARGIN_L0 = STYLE["bullet"]["margin_l0"]
    BULLET_MARGIN_L1 = STYLE["bullet"]["margin_l1"]
    _tb = STYLE["title_slide_textboxes"]
    TITLE_ISSUE_LEFT = _tb["issue"]["left"]
    TITLE_ISSUE_TOP = _tb["issue"]["top"]
    TITLE_ISSUE_WIDTH = _tb["issue"]["width"]
    TITLE_ISSUE_HEIGHT = _tb["issue"]["height"]
    TITLE_DATE_LEFT = _tb["date"]["left"]
    TITLE_DATE_TOP = _tb["date"]["top"]
    TITLE_DATE_WIDTH = _tb["date"]["width"]
    TITLE_DATE_HEIGHT = _tb["date"]["height"]
    TITLE_TB_FONT = _tb.get("font", "Zen Kaku Gothic New")
    TITLE_TB_SIZE_PT = _tb.get("size_pt", 12)
    TITLE_TB_BOLD = _tb.get("bold", True)
    TITLE_TB_COLOR = _tb.get("color_rgb", "FFFFFF")
    TITLE_TB_LINE_SPACING = _tb.get("line_spacing", 1.25)
    _tbl = STYLE["table_position"]
    TABLE_LEFT = _tbl["left"]
    TABLE_TOP = _tbl["top"]
    TABLE_WIDTH = _tbl["width"]
    TABLE_HEIGHT = _tbl["height"]
    _img = STYLE["image_layout"]
    BODY_NARROW_WIDTH = _img["body_narrow_width"]
    IMG_RIGHT_LEFT = _img["img_right_left"]
    IMG_RIGHT_TOP = _img["img_right_top"]
    IMG_RIGHT_MAX_WIDTH = _img["img_right_max_width"]
    IMG_RIGHT_MAX_HEIGHT = _img["img_right_max_height"]
    URL_HYPERLINK_COLOR = STYLE.get("url_hyperlink_color", "0563C1")
    URL_SEPARATOR_COLOR = STYLE.get("url_separator_color", "808080")
    BODY_BOLD = STYLE.get("body_bold", False)


def list_styles():
    """Return a list of available style names."""
    if not os.path.isdir(STYLES_DIR):
        return []
    return sorted(
        d for d in os.listdir(STYLES_DIR)
        if os.path.isdir(os.path.join(STYLES_DIR, d))
        and os.path.exists(os.path.join(STYLES_DIR, d, "style_config.json"))
    )


def load_template(src_path=None):
    """Load a PPTX file as template and return the Presentation object.

    Args:
        src_path: path to a PPTX file. If None, uses the bundled template.
    """
    if src_path is None:
        src_path = TEMPLATE_PATH
    return Presentation(src_path)


def analyze_template(prs):
    """Analyze slide layouts and return a summary of available layouts and their placeholders."""
    info = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "num_slides": len(prs.slides),
        "layouts": [],
    }
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {"index": i, "name": layout.name, "placeholders": []}
        for ph in layout.placeholders:
            layout_info["placeholders"].append({
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
            })
        info["layouts"].append(layout_info)
    return info


def analyze_slides(prs, max_slides=5):
    """Analyze existing slides for font, alignment, and positioning details."""
    slides_info = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "index": slide_idx + 1,
            "layout": slide.slide_layout.name,
            "shapes": [],
        }
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text[:100]
                shape_info["paragraphs"] = []
                for p in shape.text_frame.paragraphs:
                    p_info = {"text": p.text[:80], "alignment": str(p.alignment)}
                    if p.runs:
                        r = p.runs[0]
                        p_info["font_name"] = r.font.name
                        p_info["font_size"] = r.font.size
                        p_info["font_bold"] = r.font.bold
                        try:
                            p_info["font_color"] = str(r.font.color.rgb)
                        except (AttributeError, TypeError):
                            p_info["font_color"] = "inherited"
                    shape_info["paragraphs"].append(p_info)
                if shape.is_placeholder:
                    shape_info["placeholder_idx"] = shape.placeholder_format.idx
            slide_info["shapes"].append(shape_info)
        slides_info.append(slide_info)
        if slide_idx + 1 >= max_slides:
            break
    return slides_info


def get_layout(prs, name):
    """Get a slide layout by name. Returns the first matching layout."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout '{name}' not found")


def delete_all_slides(prs):
    """Delete all slides from the presentation."""
    prs_elem = prs.part._element
    sldIdLst = prs_elem.find(f"{{{NS_P}}}sldIdLst")
    for sldId in list(sldIdLst):
        rId = sldId.get(f"{{{NS_R}}}id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_slide(prs, layout):
    """Add a new slide with the given layout."""
    return prs.slides.add_slide(layout)


def set_text_in_placeholder(slide, idx, text, alignment=None):
    """Set text in a placeholder by index, optionally with alignment."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            if alignment is not None:
                for p in shape.text_frame.paragraphs:
                    p.alignment = alignment
            return shape
    return None


def _make_bullet_pPr(level=0):
    """Create XML element for a bullet paragraph property."""
    pPr = etree.SubElement(etree.Element("dummy"), f"{{{NS_A}}}pPr")
    if level == 0:
        pPr.set("marL", BULLET_MARGIN_L0)
    else:
        pPr.set("marL", BULLET_MARGIN_L1)
        pPr.set("lvl", str(level))
    etree.SubElement(pPr, f"{{{NS_A}}}buSzPct").set("val", BULLET_SIZE_PCT)
    buFont = etree.SubElement(pPr, f"{{{NS_A}}}buFont")
    buFont.set("typeface", BULLET_FONT)
    buFont.set("pitchFamily", BULLET_FONT_PITCH_FAMILY)
    buFont.set("charset", BULLET_FONT_CHARSET)
    etree.SubElement(pPr, f"{{{NS_A}}}buChar").set("char", BULLET_CHAR)
    return pPr


def _make_heading_pPr():
    """Create XML element for a non-bullet heading paragraph property."""
    pPr = etree.SubElement(etree.Element("dummy"), f"{{{NS_A}}}pPr")
    pPr.set("marL", "0")
    pPr.set("lvl", "0")
    pPr.set("indent", "0")
    pPr.set("algn", "l")
    pPr.set("rtl", "0")
    spcBef = etree.SubElement(pPr, f"{{{NS_A}}}spcBef")
    etree.SubElement(spcBef, f"{{{NS_A}}}spcPts").set("val", "0")
    spcAft = etree.SubElement(pPr, f"{{{NS_A}}}spcAft")
    etree.SubElement(spcAft, f"{{{NS_A}}}spcPts").set("val", "0")
    etree.SubElement(pPr, f"{{{NS_A}}}buNone")
    return pPr


def _make_run(text, sz=None, lang="ja"):
    """Create an XML run element with text and formatting."""
    r = etree.SubElement(etree.Element("dummy"), f"{{{NS_A}}}r")
    rPr = etree.SubElement(r, f"{{{NS_A}}}rPr")
    rPr.set("lang", lang)
    if sz:
        rPr.set("sz", str(sz))
    if BODY_BOLD:
        rPr.set("b", "1")
    rPr.set("dirty", "0")
    t = etree.SubElement(r, f"{{{NS_A}}}t")
    t.text = text
    return r


def _make_end_para_rPr(sz=None):
    """Create an endParaRPr element."""
    endPr = etree.SubElement(etree.Element("dummy"), f"{{{NS_A}}}endParaRPr")
    if sz:
        endPr.set("sz", str(sz))
    endPr.set("dirty", "0")
    return endPr


def _build_paragraph(text, para_type="bullet", level=0, font_sz=None):
    """Build an XML paragraph element with proper formatting.

    Args:
        text: paragraph text
        para_type: "bullet", "heading", or "empty"
        level: indentation level (0 or 1) for bullets
        font_sz: font size in hundredths of a point. If None, uses config body size.
    """
    if font_sz is None:
        font_sz = STYLE["fonts"]["body_size_pt100"]

    p = etree.Element(f"{{{NS_A}}}p")

    if para_type == "heading":
        pPr = _make_heading_pPr()
    elif para_type == "bullet":
        pPr = _make_bullet_pPr(level=level)
    elif para_type == "empty":
        pPr = _make_heading_pPr()
    else:
        pPr = _make_heading_pPr()

    p.append(pPr)

    if text:
        run = _make_run(text, sz=font_sz)
        p.append(run)

    p.append(_make_end_para_rPr(sz=font_sz))
    return p


def set_body_content(slide, idx, content_items):
    """Set structured body content in a placeholder with proper formatting.

    Args:
        slide: the slide object
        idx: placeholder index (typically 1 for body)
        content_items: list of dicts, each with:
            - "text": string content
            - "type": "bullet" | "heading" | "sub_bullet" | "empty"

    Example:
        content_items = [
            {"text": "Overview", "type": "heading"},
            {"text": "A bullet point", "type": "bullet"},
            {"text": "A sub-bullet", "type": "sub_bullet"},
            {"text": "", "type": "empty"},
            {"text": "Next section", "type": "heading"},
            {"text": "Another bullet", "type": "bullet"},
        ]
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            txBody = tf._txBody

            for p_elem in list(txBody.findall(f"{{{NS_A}}}p")):
                txBody.remove(p_elem)

            for item in content_items:
                text = item.get("text", "")
                item_type = item.get("type", "bullet")

                if item_type == "sub_bullet":
                    p = _build_paragraph(text, "bullet", level=1)
                elif item_type == "empty":
                    p = _build_paragraph("", "empty")
                elif item_type == "heading":
                    p = _build_paragraph(text, "heading")
                else:
                    p = _build_paragraph(text, "bullet", level=0)

                txBody.append(p)

            return shape
    return None


def set_body_bullets(slide, idx, lines, font_size=None):
    """Set body text as simple bullet points (legacy compatibility).

    For proper formatting with headings and sub-bullets, use set_body_content() instead.
    """
    items = []
    for line in lines:
        if line == "":
            items.append({"text": "", "type": "empty"})
        else:
            items.append({"text": line, "type": "bullet"})
    return set_body_content(slide, idx, items)


def set_url_in_placeholder(slide, idx, url):
    """Set URL text in a placeholder (non-clickable, legacy).

    Deprecated: Use set_url_hyperlinks() for blue clickable hyperlinks.
    """
    url_sz = str(STYLE["fonts"]["url_size_pt100"])
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            txBody = tf._txBody

            for p_elem in list(txBody.findall(f"{{{NS_A}}}p")):
                txBody.remove(p_elem)

            p = etree.SubElement(txBody, f"{{{NS_A}}}p")
            pPr = _make_heading_pPr()
            p.append(pPr)

            r = etree.SubElement(p, f"{{{NS_A}}}r")
            rPr = etree.SubElement(r, f"{{{NS_A}}}rPr")
            rPr.set("lang", "ja")
            rPr.set("sz", url_sz)
            rPr.set("dirty", "0")
            solidFill = etree.SubElement(rPr, f"{{{NS_A}}}solidFill")
            etree.SubElement(solidFill, f"{{{NS_A}}}schemeClr").set("val", "dk1")
            t = etree.SubElement(r, f"{{{NS_A}}}t")
            t.text = url

            r2 = etree.SubElement(p, f"{{{NS_A}}}r")
            rPr2 = etree.SubElement(r2, f"{{{NS_A}}}rPr")
            rPr2.set("sz", url_sz)
            rPr2.set("dirty", "0")
            t2 = etree.SubElement(r2, f"{{{NS_A}}}t")
            t2.text = " "

            p.append(_make_end_para_rPr(sz=int(url_sz)))
            return shape
    return None


def set_url_hyperlinks(slide, idx, links):
    """Set blue clickable hyperlinks in a placeholder.

    Args:
        slide: slide object
        idx: placeholder index (typically 2 for URL)
        links: list of (display_text, url) tuples
    """
    url_sz = str(STYLE["fonts"]["url_size_pt100"])
    link_color = URL_HYPERLINK_COLOR
    sep_color = URL_SEPARATOR_COLOR

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            txBody = tf._txBody

            for p_elem in list(txBody.findall(f"{{{NS_A}}}p")):
                txBody.remove(p_elem)

            p = etree.SubElement(txBody, f"{{{NS_A}}}p")
            pPr = etree.SubElement(p, f"{{{NS_A}}}pPr")
            pPr.set("marL", "0")
            pPr.set("indent", "0")
            etree.SubElement(pPr, f"{{{NS_A}}}buNone")

            for i, (text, url) in enumerate(links):
                if i > 0:
                    sep = etree.SubElement(p, f"{{{NS_A}}}r")
                    sep_rPr = etree.SubElement(sep, f"{{{NS_A}}}rPr")
                    sep_rPr.set("lang", "ja")
                    sep_rPr.set("sz", url_sz)
                    sep_rPr.set("dirty", "0")
                    sep_fill = etree.SubElement(sep_rPr, f"{{{NS_A}}}solidFill")
                    etree.SubElement(sep_fill, f"{{{NS_A}}}srgbClr").set("val", sep_color)
                    sep_t = etree.SubElement(sep, f"{{{NS_A}}}t")
                    sep_t.text = " | "

                r = etree.SubElement(p, f"{{{NS_A}}}r")
                rPr = etree.SubElement(r, f"{{{NS_A}}}rPr")
                rPr.set("lang", "ja")
                rPr.set("sz", url_sz)
                rPr.set("dirty", "0")
                fill = etree.SubElement(rPr, f"{{{NS_A}}}solidFill")
                etree.SubElement(fill, f"{{{NS_A}}}srgbClr").set("val", link_color)
                hlinkClick = etree.SubElement(rPr, f"{{{NS_A}}}hlinkClick")
                rel = slide.part.relate_to(
                    url,
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                    is_external=True,
                )
                hlinkClick.set(f"{{{NS_R}}}id", rel)
                t = etree.SubElement(r, f"{{{NS_A}}}t")
                t.text = text

            endPr = etree.SubElement(p, f"{{{NS_A}}}endParaRPr")
            endPr.set("sz", url_sz)
            endPr.set("dirty", "0")
            return shape
    return None


def remove_placeholder(slide, idx):
    """Remove a placeholder shape from a slide by index."""
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx == idx:
            sp = shape._element
            sp.getparent().remove(sp)
            return


def add_textbox(slide, left, top, width, height, text,
                font_name=None, font_size=None,
                bold=None, color=None,
                alignment=PP_ALIGN.CENTER, line_spacing=None):
    """Add a text box to a slide with specified formatting.

    Default values come from style config title textbox settings.
    """
    if font_name is None:
        font_name = TITLE_TB_FONT
    if font_size is None:
        font_size = Pt(TITLE_TB_SIZE_PT)
    if bold is None:
        bold = TITLE_TB_BOLD
    if color is None:
        c = TITLE_TB_COLOR
        color = RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    if line_spacing is None:
        line_spacing = TITLE_TB_LINE_SPACING

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.line_spacing = line_spacing
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_content_slide(prs, layout, title, body_items=None, url=None,
                      bullets=None, source_links=None):
    """Add a content slide with title, structured body content, and optional source links.

    Args:
        prs: Presentation object
        layout: slide layout object
        title: slide title text
        body_items: list of dicts for set_body_content() (preferred)
        source_links: list of (display_text, url) tuples for blue clickable hyperlinks
        url: optional plain URL string for URL placeholder (legacy, not clickable)
        bullets: legacy list of strings (uses set_body_bullets)
    """
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, PH_TITLE_IDX, title)

    if body_items:
        set_body_content(slide, PH_BODY_IDX, body_items)
    elif bullets:
        set_body_bullets(slide, PH_BODY_IDX, bullets)

    if source_links:
        set_url_hyperlinks(slide, PH_URL_IDX, source_links)
    elif url:
        set_url_in_placeholder(slide, PH_URL_IDX, url)

    return slide


def add_section_header(prs, layout, title):
    """Add a section header slide (title centered)."""
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, PH_TITLE_IDX, title, alignment=PP_ALIGN.CENTER)
    return slide


def add_title_slide(prs, layout, main_title, issue_number, date_str):
    """Add a title slide with issue number and date textboxes.

    Args:
        prs: Presentation object
        layout: title layout object
        main_title: e.g. "Weekly #dev_urandom"
        issue_number: int, displayed as "#N"
        date_str: formatted date string e.g. "2026.3.7"
    """
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, PH_TITLE_IDX, main_title, alignment=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        left=TITLE_ISSUE_LEFT, top=TITLE_ISSUE_TOP,
        width=TITLE_ISSUE_WIDTH, height=TITLE_ISSUE_HEIGHT,
        text=f"#{issue_number}",
    )

    add_textbox(
        slide,
        left=TITLE_DATE_LEFT, top=TITLE_DATE_TOP,
        width=TITLE_DATE_WIDTH, height=TITLE_DATE_HEIGHT,
        text=date_str,
        alignment=PP_ALIGN.RIGHT,
    )

    return slide


def add_next_meeting_slide(prs, layout, meeting_date_str):
    """Add a 'next meeting' slide with larger font.

    The URL placeholder is automatically removed.

    Args:
        prs: Presentation object
        layout: content layout object
        meeting_date_str: e.g. "3/14(Fri) 10:00~"
    """
    next_meeting_sz = STYLE["fonts"]["next_meeting_size_pt100"]
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, PH_TITLE_IDX, "Next Meeting")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == PH_BODY_IDX:
            tf = shape.text_frame
            txBody = tf._txBody

            for p_elem in list(txBody.findall(f"{{{NS_A}}}p")):
                txBody.remove(p_elem)

            p = _build_paragraph(meeting_date_str, "bullet", level=0, font_sz=next_meeting_sz)
            txBody.append(p)
            break

    remove_placeholder(slide, PH_URL_IDX)
    return slide


def add_table_slide(prs, layout, title, headers, rows, url=None):
    """Add a slide with a data table."""
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, PH_TITLE_IDX, title)

    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Emu(TABLE_LEFT), Emu(TABLE_TOP),
        Emu(TABLE_WIDTH), Emu(TABLE_HEIGHT),
    )
    table = table_shape.table

    tbl_pt = STYLE["fonts"]["table_size_pt100"] // 100
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(tbl_pt)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(tbl_pt)

    if url:
        set_url_in_placeholder(slide, PH_URL_IDX, url)

    return slide


def extract_issue_number(prs):
    """Extract the issue number from the title slide of the template."""
    if len(prs.slides) == 0:
        return None
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            match = re.match(r"^#(\d+)$", text)
            if match:
                return int(match.group(1))
    return None


def next_friday(from_date=None):
    """Calculate the next Friday from the given date (or today)."""
    if from_date is None:
        from_date = datetime.now()
    days_ahead = 4 - from_date.weekday()
    if days_ahead <= 0:
        days_ahead += 7
    return from_date + timedelta(days=days_ahead)


def format_presentation_date(dt):
    """Format a date as 'YYYY.M.D' for the presentation title."""
    return f"{dt.year}.{dt.month}.{dt.day}"


def format_next_meeting(dt):
    """Format the next meeting date as 'M/D(Day) 10:00~'."""
    weekdays = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    wd = weekdays[dt.weekday()]
    return f"{dt.month}/{dt.day}({wd}) 10:00~"


def read_docx(path):
    """Read a docx file and return its full text content (plain text only).

    Note: This does NOT extract hyperlinks. Use read_docx_with_hyperlinks()
    if you need actual URLs embedded in the document.
    """
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        lines.append(para.text)
    return "\n".join(lines)


def read_docx_with_hyperlinks(path):
    """Read a docx file and return text content plus all hyperlinks.

    Returns:
        (text, hyperlinks) where:
        - text: full plain text (same as read_docx)
        - hyperlinks: dict mapping display text to URL
    """
    NS_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    NS_WR = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    doc = Document(path)

    text_lines = []
    for para in doc.paragraphs:
        text_lines.append(para.text)
    text = "\n".join(text_lines)

    rels = doc.part.rels
    hyperlink_rels = {}
    for rel_id, rel in rels.items():
        if "hyperlink" in str(rel.reltype).lower():
            hyperlink_rels[rel_id] = rel.target_ref

    hyperlinks = {}
    body = doc.element.body
    for para in body.iter(f"{{{NS_W}}}p"):
        for hl in para.iter(f"{{{NS_W}}}hyperlink"):
            r_id = hl.get(f"{{{NS_WR}}}id")
            parts = []
            for run in hl.iter(f"{{{NS_W}}}t"):
                if run.text:
                    parts.append(run.text)
            display = "".join(parts).strip()
            url = hyperlink_rels.get(r_id)
            if display and url:
                hyperlinks[display] = url

    return text, hyperlinks


def read_pptx_content(path):
    """Read a PPTX file and return a text summary of each slide."""
    prs = Presentation(path)
    result = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text_frame.text)
            elif hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells]
                    slide_texts.append(" | ".join(cells))
        result.append({"slide": i + 1, "layout": slide.slide_layout.name,
                        "content": "\n".join(slide_texts)})
    return result


def verify_output(path):
    """Verify the generated PPTX and print a summary."""
    slides = read_pptx_content(path)
    print(f"Total slides: {len(slides)}")
    for s in slides:
        preview = s["content"][:100].replace("\n", " | ")
        print(f"  Slide {s['slide']} ({s['layout']}): {preview}")


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def _ensure_img_cache():
    """Create the image cache directory if it doesn't exist."""
    os.makedirs(IMG_TEMP_DIR, exist_ok=True)
    return IMG_TEMP_DIR


def fetch_og_image(page_url):
    """Extract the Open Graph image URL from a web page.

    Returns the og:image URL string, or None if not found.
    """
    try:
        resp = requests.get(
            page_url, timeout=15,
            headers={"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"},
            allow_redirects=True,
        )
        resp.raise_for_status()
        m = re.search(
            r'<meta[^>]+property=["\']og:image["\'][^>]+content=["\'](.*?)["\']',
            resp.text, re.IGNORECASE,
        )
        if not m:
            m = re.search(
                r'<meta[^>]+content=["\'](.*?)["\'][^>]+property=["\']og:image["\']',
                resp.text, re.IGNORECASE,
            )
        if m:
            return html_mod.unescape(m.group(1))
    except Exception as e:
        print(f"  [og:image] failed for {page_url}: {e}")
    return None


def download_image(url, filename=None):
    """Download an image from a URL and save it to the image cache.

    Returns:
        Absolute path to the downloaded file, or None on failure.
    """
    cache = _ensure_img_cache()
    try:
        resp = requests.get(
            url, timeout=20,
            headers={"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"},
            allow_redirects=True,
        )
        resp.raise_for_status()
        ct = resp.headers.get("content-type", "")
        if "png" in ct:
            ext = ".png"
        elif "webp" in ct:
            ext = ".webp"
        elif "gif" in ct:
            ext = ".gif"
        else:
            ext = ".jpg"
        if filename is None:
            filename = f"img_{hash(url) % 100000:05d}"
        path = os.path.join(cache, filename + ext)
        with open(path, "wb") as f:
            f.write(resp.content)
        size_kb = len(resp.content) / 1024
        print(f"  [download] {path} ({size_kb:.0f} KB)")
        return path
    except Exception as e:
        print(f"  [download] failed for {url}: {e}")
        return None


def fetch_and_download_og_image(page_url, filename=None):
    """Fetch og:image from a page URL and download it.

    Returns:
        Absolute path to the downloaded image, or None if not available.
    """
    og_url = fetch_og_image(page_url)
    if og_url:
        return download_image(og_url, filename=filename)
    return None


def resize_placeholder(slide, idx, width=None, height=None):
    """Resize a placeholder shape on a slide.

    Always writes all four geometry values explicitly to prevent python-pptx
    from zeroing out inherited values.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            orig_left = shape.left
            orig_top = shape.top
            orig_w = shape.width
            orig_h = shape.height
            shape.left = orig_left
            shape.top = orig_top
            shape.width = width if width is not None else orig_w
            shape.height = height if height is not None else orig_h
            return shape
    return None


def add_image_to_slide(slide, image_path, left, top, max_width=None, max_height=None):
    """Add an image to a slide, fitting within max dimensions while preserving aspect ratio."""
    from PIL import Image as PILImage
    try:
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size
    except Exception:
        img_w, img_h = 800, 600

    emu_per_px = 914400 / 96
    native_w = int(img_w * emu_per_px)
    native_h = int(img_h * emu_per_px)

    w, h = native_w, native_h
    if max_width and w > max_width:
        scale = max_width / w
        w = max_width
        h = int(h * scale)
    if max_height and h > max_height:
        scale = max_height / h
        h = max_height
        w = int(w * scale)

    pic = slide.shapes.add_picture(image_path, Emu(left), Emu(top), Emu(w), Emu(h))
    return pic


def add_image_right(slide, image_path):
    """Add an image to the right side of a content slide and narrow the body placeholder."""
    resize_placeholder(slide, PH_BODY_IDX, width=BODY_NARROW_WIDTH)
    pic = add_image_to_slide(
        slide, image_path,
        left=IMG_RIGHT_LEFT, top=IMG_RIGHT_TOP,
        max_width=IMG_RIGHT_MAX_WIDTH, max_height=IMG_RIGHT_MAX_HEIGHT,
    )
    if pic and pic.height < IMG_RIGHT_MAX_HEIGHT:
        offset = (IMG_RIGHT_MAX_HEIGHT - pic.height) // 2
        pic.top = Emu(IMG_RIGHT_TOP + offset)
    return pic


def add_content_slide_with_image(prs, layout, title, image_path,
                                  body_items=None, source_links=None):
    """Add a content slide with text on the left and an image on the right."""
    slide = add_content_slide(prs, layout, title,
                              body_items=body_items, source_links=source_links)
    if image_path and os.path.exists(image_path):
        add_image_right(slide, image_path)
    return slide


def cleanup_img_cache():
    """Remove all files in the image cache directory."""
    if os.path.isdir(IMG_TEMP_DIR):
        for f in os.listdir(IMG_TEMP_DIR):
            fp = os.path.join(IMG_TEMP_DIR, f)
            if os.path.isfile(fp):
                os.remove(fp)
        print(f"Cleaned up {IMG_TEMP_DIR}")
