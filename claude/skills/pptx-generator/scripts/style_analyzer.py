import os
import sys
import json
import shutil
from pptx import Presentation
from pptx.util import Pt, Emu

SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
SKILL_DIR = os.path.dirname(SCRIPTS_DIR)
STYLES_DIR = os.path.join(SKILL_DIR, "styles")
ACTIVE_STYLE_PATH = os.path.join(SKILL_DIR, "active_style.json")


def analyze_pptx(pptx_path):
    """Analyze a PPTX file and extract style information.

    Returns a dict with layout info, placeholder info, font details,
    positioning data, and color information extracted from the slides.
    """
    prs = Presentation(pptx_path)

    result = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "layouts": [],
        "slides": [],
    }

    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholders": [],
        }
        for ph in layout.placeholders:
            ph_info = {
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height,
            }
            layout_info["placeholders"].append(ph_info)
        result["layouts"].append(layout_info)

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "index": slide_idx + 1,
            "layout_name": slide.slide_layout.name,
            "shapes": [],
        }
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "is_placeholder": shape.is_placeholder,
            }
            if shape.is_placeholder:
                shape_info["placeholder_idx"] = shape.placeholder_format.idx

            if shape.has_text_frame:
                shape_info["has_text"] = True
                shape_info["text_preview"] = shape.text_frame.text[:200]
                shape_info["paragraphs"] = []
                for p in shape.text_frame.paragraphs:
                    p_info = {
                        "text": p.text[:100],
                        "alignment": str(p.alignment),
                        "level": p.level,
                    }
                    if p.runs:
                        r = p.runs[0]
                        p_info["font_name"] = r.font.name
                        p_info["font_size"] = r.font.size
                        p_info["font_bold"] = r.font.bold
                        try:
                            p_info["font_color_rgb"] = str(r.font.color.rgb)
                        except (AttributeError, TypeError):
                            p_info["font_color_rgb"] = None
                    shape_info["paragraphs"].append(p_info)

            if hasattr(shape, "has_table") and shape.has_table:
                shape_info["has_table"] = True
                shape_info["table_rows"] = len(shape.table.rows)
                shape_info["table_cols"] = len(shape.table.columns)

            slide_info["shapes"].append(shape_info)
        result["slides"].append(slide_info)

    return result


def _find_layout_by_role(analysis):
    """Identify which layouts serve as title/section and content layouts.

    Heuristic:
    - Title layout: exactly 1 placeholder (title only).
    - Content layout: the layout with the most placeholders (2+).
      More placeholders typically means title + body + url, which is the
      richest content layout available.
    """
    title_layout = None
    content_layout = None
    max_ph_count = 0

    for layout in analysis["layouts"]:
        phs = layout["placeholders"]
        if len(phs) == 1 and phs[0]["idx"] == 0:
            if title_layout is None:
                title_layout = layout
        elif len(phs) >= 2 and len(phs) > max_ph_count:
            content_layout = layout
            max_ph_count = len(phs)

    if title_layout is None and analysis["layouts"]:
        title_layout = analysis["layouts"][0]
    if content_layout is None and len(analysis["layouts"]) > 1:
        content_layout = analysis["layouts"][1]

    return title_layout, content_layout


def _extract_font_info(analysis):
    """Extract dominant font information from slide content."""
    font_names = {}
    font_sizes = {}
    font_bolds = {}

    for slide in analysis["slides"]:
        for shape in slide["shapes"]:
            if "paragraphs" not in shape:
                continue
            for p in shape["paragraphs"]:
                fn = p.get("font_name")
                fs = p.get("font_size")
                fb = p.get("font_bold")
                if fn:
                    font_names[fn] = font_names.get(fn, 0) + 1
                if fs:
                    font_sizes[fs] = font_sizes.get(fs, 0) + 1
                if fb is not None:
                    font_bolds[fb] = font_bolds.get(fb, 0) + 1

    dominant_font = max(font_names, key=font_names.get) if font_names else "Noto Sans JP"
    dominant_size = max(font_sizes, key=font_sizes.get) if font_sizes else 165100

    return {
        "dominant_font": dominant_font,
        "dominant_size_emu": dominant_size,
        "font_frequency": font_names,
        "size_frequency": {str(k): v for k, v in font_sizes.items()},
    }


def _extract_textbox_info(analysis):
    """Extract textbox positioning from slides (non-placeholder shapes with text)."""
    textboxes = []
    for slide in analysis["slides"]:
        for shape in slide["shapes"]:
            if not shape.get("is_placeholder") and shape.get("has_text"):
                textboxes.append({
                    "text_preview": shape.get("text_preview", "")[:50],
                    "left": shape["left"],
                    "top": shape["top"],
                    "width": shape["width"],
                    "height": shape["height"],
                    "slide_index": slide["index"],
                    "paragraphs": shape.get("paragraphs", []),
                })
    return textboxes


def _build_available_layouts(analysis):
    """Build available_layouts list from PPTX analysis data."""
    layouts = []
    for layout in analysis.get("layouts", []):
        entry = {
            "name": layout["name"],
            "index": layout["index"],
            "placeholders": [
                {"idx": p["idx"], "name": p["name"]}
                for p in layout["placeholders"]
            ],
        }
        layouts.append(entry)
    return layouts


def build_style_config(analysis, style_name="custom"):
    """Build a style_config.json dict from PPTX analysis.

    This produces a best-effort configuration. The user/agent should review
    and adjust values as needed after generation.
    """
    title_layout, content_layout = _find_layout_by_role(analysis)
    font_info = _extract_font_info(analysis)
    textboxes = _extract_textbox_info(analysis)

    body_size_pt100 = int(font_info["dominant_size_emu"] / 12700) if font_info["dominant_size_emu"] else 1300

    ph_indices = {"title_idx": 0, "body_idx": 1, "url_idx": 2}
    if content_layout:
        phs = sorted(content_layout["placeholders"], key=lambda x: x["idx"])
        if len(phs) >= 1:
            ph_indices["title_idx"] = phs[0]["idx"]
        if len(phs) >= 2:
            ph_indices["body_idx"] = phs[1]["idx"]
        if len(phs) >= 3:
            ph_indices["url_idx"] = phs[2]["idx"]

    title_tb = {
        "issue": {"left": 1607335, "top": 3756390, "width": 5973000, "height": 369300},
        "date": {"left": 2820600, "top": 450625, "width": 5973000, "height": 369300},
        "font": font_info["dominant_font"],
        "size_pt": 12,
        "bold": True,
        "color_rgb": "FFFFFF",
        "line_spacing": 1.25,
    }
    if textboxes:
        tb = textboxes[0]
        title_tb["issue"]["left"] = tb["left"]
        title_tb["issue"]["top"] = tb["top"]
        title_tb["issue"]["width"] = tb["width"]
        title_tb["issue"]["height"] = tb["height"]
        if tb.get("paragraphs"):
            p = tb["paragraphs"][0]
            if p.get("font_name"):
                title_tb["font"] = p["font_name"]
            if p.get("font_size"):
                title_tb["size_pt"] = int(p["font_size"] / 12700)
            if p.get("font_bold") is not None:
                title_tb["bold"] = p["font_bold"]
            if p.get("font_color_rgb"):
                title_tb["color_rgb"] = p["font_color_rgb"]
        if len(textboxes) > 1:
            tb2 = textboxes[1]
            title_tb["date"]["left"] = tb2["left"]
            title_tb["date"]["top"] = tb2["top"]
            title_tb["date"]["width"] = tb2["width"]
            title_tb["date"]["height"] = tb2["height"]

    slide_w = analysis["slide_width"]
    body_narrow = int(slide_w * 0.5)
    img_left = int(slide_w * 0.6)
    img_max_w = int(slide_w * 0.35)

    config = {
        "style_name": style_name,
        "slide_dimensions": {
            "width_emu": analysis["slide_width"],
            "height_emu": analysis["slide_height"],
        },
        "layouts": {
            "title": {
                "name": title_layout["name"] if title_layout else "Title",
                "description": "Title slide / section headers",
            },
            "content": {
                "name": content_layout["name"] if content_layout else "Content",
                "description": "Content slides",
            },
        },
        "placeholders": ph_indices,
        "fonts": {
            "master_default": font_info["dominant_font"],
            "title_textbox_font": title_tb["font"],
            "body_size_pt100": body_size_pt100,
            "url_size_pt100": 1000,
            "table_size_pt100": 1000,
            "next_meeting_size_pt100": 1400,
            "title_textbox_size_pt": title_tb["size_pt"],
        },
        "bullet": {
            "char": "l",
            "font": "Wingdings",
            "pitch_family": "2",
            "charset": "2",
            "size_pct": "75000",
            "margin_l0": "431800",
            "margin_l1": "889000",
        },
        "title_slide_textboxes": title_tb,
        "table_position": {
            "left": 650888,
            "top": 1136755,
            "width": 7842200,
            "height": 3444000,
        },
        "image_layout": {
            "body_narrow_width": body_narrow,
            "img_right_left": img_left,
            "img_right_top": 1228013,
            "img_right_max_width": img_max_w,
            "img_right_max_height": 3015900,
        },
        "url_hyperlink_color": "0563C1",
        "url_separator_color": "808080",
        "body_bold": False,
        "available_layouts": _build_available_layouts(analysis),
        "content_guidelines": {
            "lines_per_slide_min": 5,
            "lines_per_slide_max": 13,
            "bold_detail": "本文テキストには太字を使用しない",
            "section_pattern": "見出し → 箇条書き → 空行のパターンで構成",
        },
    }

    return config


def update_style(pptx_path, style_name=None, set_active=True):
    """Main entry point: analyze a PPTX and create a new style in styles/<name>/.

    1. Analyzes the provided PPTX
    2. Creates styles/<style_name>/ directory
    3. Writes style_config.json into it (available_layouts, content_guidelines 含む)
    4. Copies the PPTX as template.pptx into it
    5. Optionally updates active_style.json to point to this style

    Args:
        pptx_path: path to the PPTX file to use as the new style source
        style_name: name for this style (defaults to filename without extension)
        set_active: if True, update active_style.json to this style

    Returns:
        (config, analysis) tuple
    """
    if style_name is None:
        style_name = os.path.splitext(os.path.basename(pptx_path))[0]

    style_dir = os.path.join(STYLES_DIR, style_name)
    config_path = os.path.join(style_dir, "style_config.json")
    template_path = os.path.join(style_dir, "template.pptx")

    os.makedirs(style_dir, exist_ok=True)

    print(f"Analyzing: {pptx_path}")
    analysis = analyze_pptx(pptx_path)

    print(f"Building style config for '{style_name}'...")
    config = build_style_config(analysis, style_name=style_name)

    print(f"Writing: {config_path}")
    with open(config_path, "w", encoding="utf-8") as f:
        json.dump(config, f, indent=2, ensure_ascii=False)

    print(f"Copying template: {pptx_path} -> {template_path}")
    shutil.copy2(pptx_path, template_path)

    if set_active:
        print(f"Setting active style: {style_name}")
        with open(ACTIVE_STYLE_PATH, "w", encoding="utf-8") as f:
            json.dump({"active": style_name}, f, indent=2, ensure_ascii=False)
            f.write("\n")

    print(f"Style '{style_name}' created at: {style_dir}")
    return config, analysis


def print_analysis_summary(analysis):
    """Print a human-readable summary of the PPTX analysis."""
    dims = analysis
    print(f"Slide size: {dims['slide_width']:,} x {dims['slide_height']:,} EMU")
    print(f"           ({dims['slide_width']/914400:.2f} x {dims['slide_height']/914400:.2f} inch)")
    print(f"Layouts: {len(dims['layouts'])}")
    for layout in dims["layouts"]:
        phs = ", ".join(f"idx={p['idx']}:{p['name']}" for p in layout["placeholders"])
        print(f"  [{layout['index']}] {layout['name']}: {phs}")
    print(f"Slides: {len(dims['slides'])}")
    for slide in dims["slides"]:
        shapes_desc = ", ".join(s["name"] for s in slide["shapes"][:5])
        print(f"  Slide {slide['index']} ({slide['layout_name']}): {shapes_desc}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python style_analyzer.py <path-to-pptx> [style-name] [--no-activate]")
        print("  Creates a new style in styles/<style-name>/ with style_config.json and template.pptx")
        print("  By default, also sets it as the active style.")
        sys.exit(1)
    path = sys.argv[1]
    name = sys.argv[2] if len(sys.argv) > 2 and not sys.argv[2].startswith("--") else None
    activate = "--no-activate" not in sys.argv
    config, analysis = update_style(path, style_name=name, set_active=activate)
    print("\n--- Analysis Summary ---")
    print_analysis_summary(analysis)
    print(f"\n--- Generated Config ---")
    print(json.dumps(config, indent=2, ensure_ascii=False))
